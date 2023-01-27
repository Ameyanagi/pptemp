from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN,MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import pandas as pd
import numpy as np
from lxml import etree
import re

import glob
from PIL import Image
import tqdm

class pptemp(object):
        
    align_dict = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    vertical_dict = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    # color_dict = {"black": RGBColor(0x00, 0x00, 0x00),
    #               "red": RGBColor(0xFF, 0x00, 0x00),
    #               "green": RGBColor(0x00, 0xFF, 0x00),
    #               "blue": RGBColor(0x00, 0x00, 0xFF),
    #               "white": RGBColor(0xFF, 0xFF, 0xFF)
    #               }
    
    color_dict = {"black": RGBColor(0x40, 0x40, 0x40),
                  "red": RGBColor(0xFF, 0x40, 0x40),
                  "green": RGBColor(0x40, 0xFF, 0x40),
                  "blue": RGBColor(0x40, 0x40, 0xFF),
                  "white": RGBColor(0xFF, 0xFF, 0xFF)
                  }
    
    
    def __init__(self, path=None, *args):
        super(pptemp, self).__init__(*args)
        self.prs = Presentation(path)
        self.blank = self.prs.slide_layouts[6]
        
        # Set slide size to 16*9 otherwise the path is specified
        if path == None:
            self.prs.slide_width = 12192000
            self.prs.slide_height = 6858000
    
    #Slides    
        
    def add_title_slide(self,title = "Title", subtitle = "Name", align = "center", vertical = "middle", 
                        font_name = "Meiryo", font_size_title = 44, font_size_subtitle = 18, font_bold = True, font_italic = False, font_underline = False, font_color = "black"):
        slide = self.prs.slides.add_slide(self.blank)
        
        # Add title textbox
        slide, textbox = self.add_textbox(slide, title, 10, 25, 80, 30, align = align, vertical = vertical, 
                    font_name = font_name, font_size = font_size_title, font_bold = font_bold, font_italic = font_italic, font_underline = font_underline, font_color = font_color)
                
        # Add subtitle textbox
        slide, textbox = self.add_textbox(slide, subtitle, 20, 60, 60, 20, align = align, vertical = vertical, 
                    font_name = font_name, font_size = font_size_subtitle, font_bold = font_bold, font_italic = font_italic, font_underline = font_underline, font_color = font_color)
                
        return slide
    
    
    def add_content_slide(self,title = "Title", use_bar = True, align = "left", vertical = "top", 
                    font_name = "Meiryo", font_size = 30, font_bold = True, font_italic = False, font_underline = False, font_color = "black"):
        
        # Create New Slide
        slide = self.prs.slides.add_slide(self.blank)
                
        # Add title textbox
        slide, textbox = self.add_textbox(slide, title, 1, 2, 95, 5, align = align, vertical = vertical, 
                    font_name = font_name, font_size = font_size, font_bold = font_bold, font_italic = font_italic, font_underline = font_underline, font_color = font_color)
        
        if use_bar == True:
            slide,_ = self.add_bar(slide)
        else:
            pass
        
        return slide

    # Contents
    
    def add_textbox(self, slide, text = "", left = 0, top = 0, width = 20, height = 5, align = "center", vertical = "middle", 
                    font_name = "Meiryo", font_size = 18, font_bold = False, font_italic = False, font_underline = False, font_color = "black"):
        # Create textbox by %
                             
        textbox = slide.shapes.add_textbox(self.prs.slide_width*left/100,
                                           self.prs.slide_height*top/100,
                                           self.prs.slide_width*width/100,
                                           self.prs.slide_height*height/100)
        
        
        # Clear textbox
        text_frame = textbox.text_frame

        # Set vertical anchor
        text_frame.vertical_anchor = self.vertical_dict[vertical]
        
        # Set word wrap
        text_frame.word_wrap = True
        
        # Get paragraph
        p = text_frame.paragraphs[0]        
        run = p.add_run()
        run.text = text
        
        # Set fonts
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = font_bold
        run.font.italic = font_italic
        run.font.underline = font_underline
        run.font.color.rgb = self.color_dict[font_color]
        
        p.alignment = self.align_dict[align]
        
        return slide, textbox
    
    def add_picture(self, slide, path, left=0, top=0, width=50, height=50):
        
        picture = slide.shapes.add_picture(*self.calc_center_img(path, left, top, width, height))
        
        return slide, picture
    
    def add_picture_label(self, slide, path, left=0, top=0, width=50, height=50, label = "", label_position = "top", label_height = 5,
                          align = "center", vertical = "middle", font_name = "Meiryo", font_size = 18, font_bold = False, 
                          font_italic = False, font_underline = True, font_color = "black"):
        
        if label_position == "top":
            picture_top = top+label_height
            picture_height = height-label_height
            label_top = top
        else:
            picture_top = top
            picture_height = height-label_height
            label_top = top+picture_height
            
        picture = slide.shapes.add_picture(*self.calc_center_img(path, left, picture_top, width, picture_height))
        textbox = self.add_textbox(slide, label, left, label_top, width, label_height, align = align, vertical = vertical, font_name = font_name, 
                                    font_size = font_size, font_bold = font_bold, font_italic = font_italic, 
                                    font_underline = font_underline, font_color = font_color)
            
        
        return slide, picture, textbox

    def calc_center_img(self, path, left, top, width, height):
        
        # Obtain resolution
        # SVG file is not supported.
        aspect_ratio_img = self.obtain_aspect_ratio(path)
                
        width_box = self.prs.slide_width*width/100
        height_box = self.prs.slide_height*height/100
        
        aspect_ratio_box = height_box/width_box
        
        
        if aspect_ratio_box > aspect_ratio_img:
            width_img = width_box
            height_img = width_box*aspect_ratio_img
        else:
            height_img = height_box
            width_img = height_img/aspect_ratio_img
                
        left_box = self.prs.slide_width*left/100+(width_box-width_img)/2
        top_box = self.prs.slide_height*top/100+(height_box-height_img)/2

        return path, left_box, top_box, width_img, height_img
    
    def obtain_aspect_ratio(self, path):
        
        file_type = path.split(".")[-1].lower()
        image_type = ["png", "jpg", "jpeg", "bmp", "gif"]
                        
        
        if file_type in image_type:
            img = Image.open(path)
            width_img, height_img = img.size
        elif file_type == "svg":
            tree = etree.parse(path, parser=etree.XMLParser(huge_tree=True))
            svg = tree.getroot()
            if 'width' not in svg.keys() or 'height' not in svg.keys():
                raise Exception('SVG header must contain width and height attributes')
            width_img = self.parse_length(svg.get('width'))
            height_img = self.parse_length(svg.get('height'))                
        else:
            pass
        
        
        return height_img/width_img
        
        
         
    def add_bar(self, slide, left = 0, top = 10, width = 100, height = 2):
        # color of the gradient
        stops_color = [RGBColor(0xE0,0xE5,0xF7), RGBColor(0x95, 0xAB, 0xEA), ]
        
        # Create textbox by %
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                     self.prs.slide_width*left/100,
                                     self.prs.slide_height*top/100,
                                     self.prs.slide_width*width/100,
                                     self.prs.slide_height*height/100)
        
        # Setting the gradient color
        fill = bar.fill
        fill.gradient()
        fill.gradient_angle = 0
        fill.gradient_stops[0].color.rgb = stops_color[0]
        fill.gradient_stops[1].color.rgb = stops_color[1]
        
        # Setting Lines
        line = bar.line
        line.width = 1
        line.color.rgb = self.color_dict["white"]
        
        # Setting Shadows
        bar.shadow.inherit = False
              
        return slide, bar
    
    def add_table_from_df(self, slide, df, left, top, width, height, colnames = None,
                          align = "center", vertical = "middle", font_name = "Meiryo", font_size = 18, 
                          font_bold = False, font_italic = False, font_underline = False, font_color = "black"):
        
        left_table = int(self.prs.slide_width*left/100)
        top_table = int(self.prs.slide_height*top/100)
        width_table = int(self.prs.slide_width*width/100)
        height_table = int(self.prs.slide_height*height/100)
        
        rows, cols = df.shape
        table = slide.shapes.add_table(rows+1, cols, left_table, top_table, width_table, height_table)

        if colnames is None:
            colnames = list(df.columns)

        # Insert the column names
        for col_index, col_name in enumerate(colnames):
            cell = table.table.cell(0,col_index)
            self.text_to_cell(cell, str(col_name), align = align, vertical = vertical, font_name = font_name, font_size = font_size, 
                     font_bold = font_bold, font_italic = font_italic, font_underline = font_underline, font_color = font_color)

            # Set bottom border
            self._set_cell_border(cell, border_width=str(Pt(1.5)), border_position="top")
            self._set_cell_border(cell, border_width=str(Pt(1)), border_position="bottom")
            
            # Set background color
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                        
        table_values = df.values
        
        # Insert the values to table
        for row in range(rows):
            for col in range(cols):
                
                cell = table.table.cell(row+1, col)
                self.text_to_cell(cell, str(table_values[row][col]), align = align, vertical = vertical, font_name = font_name, font_size = font_size, 
                     font_bold = font_bold, font_italic = font_italic, font_underline = font_underline, font_color = font_color)
                
                # Set bottom border
                if row == rows-1:
                    self._set_cell_border(cell, border_width=str(Pt(1.5)), border_position="bottom")

                # Set background color
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
        return slide, table        
    
    def text_to_cell(self, cell, text = "", align = "center", vertical = "middle", font_name = "Meiryo", font_size = 18, 
                     font_bold = False, font_italic = False, font_underline = False, font_color = "black"):
        cell.vertical_anchor = self.vertical_dict[vertical]
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = text
        
        # Set fonts
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = font_bold
        run.font.italic = font_italic
        run.font.underline = font_underline
        run.font.color.rgb = self.color_dict[font_color]
        
        p.alignment = self.align_dict[align]
    
    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element
    
    def _set_cell_border(self, cell, border_color="000000", border_width="12700", border_position = "left right top bottom"):
        """ Hack function to enable the setting of border width and border color
            - left border
            - right border
            - top border
            - bottom border
        """
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        if "left" in border_position:
            # Left Cell Border
            lnL = self.SubElement(tcPr, 'a:lnL', w=border_width, cap='flat', cmpd='sng', algn='ctr')
            lnL_solidFill = self.SubElement(lnL, 'a:solidFill')
            lnL_srgbClr = self.SubElement(lnL_solidFill, 'a:srgbClr', val=border_color)
            lnL_prstDash = self.SubElement(lnL, 'a:prstDash', val='solid')
            lnL_round_ = self.SubElement(lnL, 'a:round')
            lnL_headEnd = self.SubElement(lnL, 'a:headEnd', type='none', w='med', len='med')
            lnL_tailEnd = self.SubElement(lnL, 'a:tailEnd', type='none', w='med', len='med')

        if "right" in border_position:
            # Right Cell Border
            lnR = self.SubElement(tcPr, 'a:lnR', w=border_width, cap='flat', cmpd='sng', algn='ctr')
            lnR_solidFill = self.SubElement(lnR, 'a:solidFill')
            lnR_srgbClr = self.SubElement(lnR_solidFill, 'a:srgbClr', val=border_color)
            lnR_prstDash = self.SubElement(lnR, 'a:prstDash', val='solid')
            lnR_round_ = self.SubElement(lnR, 'a:round')
            lnR_headEnd = self.SubElement(lnR, 'a:headEnd', type='none', w='med', len='med')
            lnR_tailEnd = self.SubElement(lnR, 'a:tailEnd', type='none', w='med', len='med')

        if "top" in border_position:
            # Top Cell Border
            lnT = self.SubElement(tcPr, 'a:lnT', w=border_width, cap='flat', cmpd='sng', algn='ctr')
            lnT_solidFill = self.SubElement(lnT, 'a:solidFill')
            lnT_srgbClr = self.SubElement(lnT_solidFill, 'a:srgbClr', val=border_color)
            lnT_prstDash = self.SubElement(lnT, 'a:prstDash', val='solid')
            lnT_round_ = self.SubElement(lnT, 'a:round')
            lnT_headEnd = self.SubElement(lnT, 'a:headEnd', type='none', w='med', len='med')
            lnT_tailEnd = self.SubElement(lnT, 'a:tailEnd', type='none', w='med', len='med')

        if "bottom" in border_position:
            # Bottom Cell Border
            lnB = self.SubElement(tcPr, 'a:lnB', w=border_width, cap='flat', cmpd='sng', algn='ctr')
            lnB_solidFill = self.SubElement(lnB, 'a:solidFill')
            lnB_srgbClr = self.SubElement(lnB_solidFill, 'a:srgbClr', val=border_color)
            lnB_prstDash = self.SubElement(lnB, 'a:prstDash', val='solid')
            lnB_round_ = self.SubElement(lnB, 'a:round')
            lnB_headEnd = self.SubElement(lnB, 'a:headEnd', type='none', w='med', len='med')
            lnB_tailEnd = self.SubElement(lnB, 'a:tailEnd', type='none', w='med', len='med')
    
    # Util
    
    def parse_length(self, value, def_units='px'):
        """Parses value as SVG length and returns it in pixels, or a negative scale (-1 = 100%).
        obtained from https://github.com/Zverik/svg-resize/blob/master/svg_resize.py"""
        
        if not value:
            return 0.0
        parts = re.match(r'^\s*(-?\d+(?:\.\d+)?)\s*(px|in|cm|mm|pt|pc|%)?', value)
        if not parts:
            raise Exception('Unknown length format: "{}"'.format(value))
        num = float(parts.group(1))
        units = parts.group(2) or def_units
        if units == 'px':
            return num
        elif units == 'pt':
            return num * 1.25
        elif units == 'pc':
            return num * 15.0
        elif units == 'in':
            return num * 90.0
        elif units == 'mm':
            return num * 3.543307
        elif units == 'cm':
            return num * 35.43307
        elif units == '%':
            return -num / 100.0
        else:
            raise Exception('Unknown length units: {}'.format(units))
        
    def calc_align_img(self, path_list, left, top, width, height):
                
        columns = int(np.sqrt(len(path_list)))
        rows = np.ceil(len(path_list)/columns)
        
        width = width/rows
        height = height/columns
                
        results = []
                
        for i in range(len(path_list)):

            row = i%rows
            column = i//rows
            
            if column == len(path_list)//rows:
                residue = rows - len(path_list)%rows
            else:
                residue = 0
            
            results.append([path_list[i], left+width*(row+residue/2), top+height*column, width, height])
            
        return results
    
    def get_img_list(self, path, file_regex = re.compile(r".*[_/\\](.*)\.[a-zA-Z]+")):
        file_list = glob.glob(path)
        file_sep_list = []
        file_list.sort()
        
        for file in file_list:
            file_sep_list.append(re.findall(file_regex, file)[0])
                    
        return file_list, file_sep_list

    def add_figure_label_slide(self, dir_path = "./fig/*/", img_path = "*.png", left=0, top=12, width=100, height=88, 
                               file_regex = re.compile(r".*[_/\\](.*)\.[a-zA-Z]+"), dir_regex = re.compile(r".*[_/\\](.*)[/\\]"), use_label = True, use_bar = True, label_position = "top", title_font_size = 30, label_font_size=18):
        # Create slides from figures
        dir_list = glob.glob(dir_path)
        dir_list.sort()

        for dir in tqdm.tqdm(dir_list):
            file_list, file_sep_list = self.get_img_list(dir+img_path, file_regex=file_regex)
            
            if len(file_list) == 0:
                return
            
            name = re.findall(dir_regex, dir)
            slide_title = name[0]
            slide  = self.add_content_slide(slide_title, use_bar=use_bar, font_size=title_font_size)

            img_list = self.calc_align_img(file_list, left, top, width, height)
            
            for i in range(len(file_list)):
                if use_label == True:
                    self.add_picture_label(slide, *img_list[i], file_sep_list[i], align="center", label_position=label_position, font_size=label_font_size)
                else:
                    self.add_picture(slide, *img_list[i])
        
        return slide
        
    # Template


    def save(self, path="test.pptx"):
        self.prs.save(path)
    
    
        
    
        